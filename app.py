# app.py — Bora Alí (capitais) — Versão final pronta para GitHub / Streamlit Cloud
# Requisitos: pandas, streamlit, plotly, pydeck, prophet, python-pptx, numpy

import os
import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import pydeck as pdk
from prophet import Prophet
from prophet.plot import plot_plotly
from datetime import datetime

# -----------------------------
# PATHS — ajustar apenas se quiser (por padrão espera data/ no repo)
DATA_PATH = "data/INMET_ANAC_ROTAS_APENAS_CAPITAIS.csv"
OUTPUT_PATH = "outputs"
os.makedirs(OUTPUT_PATH, exist_ok=True)

# -----------------------------
# Branding — Azul + Coral (limpo, jovem, profissional)
PRIMARY = "#006DCE"
ACCENT = "#FF6B4A"
BG = "#F7F9FB"
TEXT = "#0F172A"

st.set_page_config(page_title="Bora Alí — Capitais", layout="wide", initial_sidebar_state="expanded")

st.markdown(f"""
<style>
body {{background-color: {BG};}}
h1 {{color: {PRIMARY}; font-weight:800;}}
h2,h3 {{color: {TEXT};}}
.stButton>button {{background-color:{PRIMARY}; color:white; border-radius:8px;}}
.card {{border-radius:10px; padding:12px; box-shadow: 0 8px 24px rgba(15,23,42,0.06); background:white}}
.small-muted {{color:#6b7280; font-size:13px}}
.legend-bora {{font-weight:600; color:{TEXT}}}
</style>
""", unsafe_allow_html=True)

# -----------------------------
# Helper: Lista completa de capitais do Brasil
CAPITAIS = {
    'Rio Branco': (-9.97499, -67.8243),
    'Maceió': (-9.649847, -35.70895),
    'Macapá': (0.034934, -51.0694),
    'Manaus': (-3.119028, -60.021731),
    'Salvador': (-12.97139, -38.50139),
    'Fortaleza': (-3.71722, -38.543366),
    'Brasília': (-15.793889, -47.882778),
    'Vitória': (-20.3155, -40.3128),
    'Goiânia': (-16.686891, -49.264788),
    'São Luís': (-2.52972, -44.30278),
    'Cuiabá': (-15.601415, -56.097892),
    'Campo Grande': (-20.4433, -54.6465),
    'Belo Horizonte': (-19.916681, -43.934493),
    'Belém': (-1.455833, -48.504444),
    'João Pessoa': (-7.119495, -34.845011),
    'Curitiba': (-25.429596, -49.271272),
    'Recife': (-8.047562, -34.8770),
    'Teresina': (-5.08921, -42.8016),
    'Rio de Janeiro': (-22.906847, -43.172896),
    'Natal': (-5.795, -35.209),
    'Porto Alegre': (-30.034647, -51.217658),
    'Porto Velho': (-8.7608, -63.9039),
    'Boa Vista': (2.8196, -60.6733),
    'Florianópolis': (-27.595377, -48.548046),
    'Aracaju': (-10.9472, -37.0731),
    'São Paulo': (-23.55052, -46.633308),
    'Palmas': (-10.184, -48.333),
    # (se faltar alguma capital no seu CSV, ajuste o nome para bater com DEST/ROTA)
}

# -----------------------------
# Load data
@st.cache_data(ttl=900)
def load_data(path):
    df = pd.read_csv(path, low_memory=False)
    # padroniza colunas esperadas
    for c in ["COMPANHIA","ROTA","DESTINO","ORIGEM","TARIFA","TEMP_MEDIA","ANO","MES"]:
        if c not in df.columns:
            df[c] = pd.NA
    # limpeza rápida
    df['ROTA'] = df['ROTA'].astype(str).str.replace(' - ', ' → ').str.strip()
    df['COMPANHIA'] = df['COMPANHIA'].astype(str).str.upper().str.strip()
    df['TARIFA'] = pd.to_numeric(df['TARIFA'], errors='coerce')
    df['TEMP_MEDIA'] = pd.to_numeric(df['TEMP_MEDIA'], errors='coerce')
    df['ANO'] = pd.to_numeric(df['ANO'], errors='coerce').fillna(0).astype(int)
    df['MES'] = pd.to_numeric(df['MES'], errors='coerce').fillna(0).astype(int)
    df['DATA'] = pd.to_datetime(df['ANO'].astype(str) + "-" + df['MES'].astype(str).str.zfill(2) + "-01", errors='coerce')
    return df

try:
    df = load_data(DATA_PATH)
except FileNotFoundError:
    st.error("Arquivo não encontrado: coloque INMET_ANAC_ROTAS_APENAS_CAPITAIS.csv em /data/ no repositório.")
    st.stop()

# -----------------------------
# Sidebar — filtros essenciais (limpo)
st.sidebar.header("Filtros — Bora Alí")
anos = sorted(df['ANO'].dropna().unique())
sel_anos = st.sidebar.multiselect("Ano", anos, default=anos if anos else [])
sel_meses = st.sidebar.multiselect("Mês", list(range(1,13)), default=list(range(1,13)))
companias = sorted(df['COMPANHIA'].dropna().unique())
sel_comp = st.sidebar.multiselect("Companhia", companias, default=companias if companias else [])
# capitais para focar (orig/dest)
cap_options = sorted(list(CAPITAIS.keys()))
sel_caps = st.sidebar.multiselect("Capitais (origem/destino)", cap_options, default=['São Paulo','Rio de Janeiro','Recife','Brasília','Manaus'])

# quick filter application
dff = df.copy()
if sel_anos:
    dff = dff[dff['ANO'].isin(sel_anos)]
if sel_meses:
    dff = dff[dff['MES'].isin(sel_meses)]
if sel_comp:
    dff = dff[dff['COMPANHIA'].isin(sel_comp)]

# parse ROTA into origin/dest (robusto)
def parse_rota(rt):
    if pd.isna(rt): return (None,None)
    if '→' in rt:
        parts = [p.strip() for p in rt.split('→') if p.strip()]
        if len(parts) >= 2: return parts[0], parts[-1]
    # fallback: try dash
    if '-' in rt:
        parts = [p.strip() for p in rt.split('-') if p.strip()]
        if len(parts) >= 2: return parts[0], parts[-1]
    return (None,None)

parsed = dff['ROTA'].apply(lambda r: pd.Series(parse_rota(r), index=['ORIG_ROTA','DEST_ROTA']))
dff = pd.concat([dff, parsed], axis=1)

# keep only routes where both origin and dest are capitals and in user selection
dff = dff[dff['ORIG_ROTA'].isin(sel_caps) & dff['DEST_ROTA'].isin(sel_caps)].copy()

# small guard
if dff.shape[0] == 0:
    st.warning("Nenhum registro após filtros. Ajuste Ano/Mês/Companhia/Capitais.")
    st.stop()

# -----------------------------
# Aggregate for maps and summaries
agg_cap = (
    dff.groupby('DEST_ROTA')
    .agg(tarifa_media=('TARIFA','mean'), temp_media=('TEMP_MEDIA','mean'), registros=('TARIFA','count'))
    .reset_index()
    .rename(columns={'DEST_ROTA':'CAPITAL'})
)
agg_cap['lat'] = agg_cap['CAPITAL'].map(lambda x: CAPITAIS.get(x,(np.nan,np.nan))[0])
agg_cap['lon'] = agg_cap['CAPITAL'].map(lambda x: CAPITAIS.get(x,(np.nan,np.nan))[1])

# prepare routes (unique)
routes = (
    dff.groupby(['ORIG_ROTA','DEST_ROTA','ROTA','COMPANHIA'])
    .agg(tarifa_media=('TARIFA','mean'), registros=('TARIFA','count'))
    .reset_index()
)
routes['olat'] = routes['ORIG_ROTA'].map(lambda x: CAPITAIS.get(x,(np.nan,np.nan))[0])
routes['olon'] = routes['ORIG_ROTA'].map(lambda x: CAPITAIS.get(x,(np.nan,np.nan))[1])
routes['dlat'] = routes['DEST_ROTA'].map(lambda x: CAPITAIS.get(x,(np.nan,np.nan))[0])
routes['dlon'] = routes['DEST_ROTA'].map(lambda x: CAPITAIS.get(x,(np.nan,np.nan))[1])
routes = routes.dropna(subset=['olat','olon','dlat','dlon']).copy()

# -----------------------------
# Layout: top metrics + maps
st.markdown("## Resumo rápido")
c1, c2, c3, c4 = st.columns([1.5,1.2,1.2,1.2])
with c1:
    st.markdown("**Registros (filtrados)**")
    st.metric("", f"{dff.shape[0]:,}")
with c2:
    st.markdown("**Tarifa média (R$)**")
    st.metric("", f"{dff['TARIFA'].mean():.0f}")
with c3:
    st.markdown("**Temperatura média (°C)**")
    st.metric("", f"{dff['TEMP_MEDIA'].mean():.1f}")
with c4:
    st.markdown("**Rotas únicas**")
    st.metric("", f"{routes['ROTA'].nunique():,}")

st.markdown("---")

left, right = st.columns([1,1])
with left:
    st.subheader("Mapa 1 — Capitais: Tarifa (tamanho) × Temperatura (cor)")
    fig_map = px.scatter_mapbox(
        agg_cap.dropna(subset=['lat','lon']),
        lat='lat', lon='lon', size='tarifa_media', color='temp_media',
        hover_name='CAPITAL', hover_data={'tarifa_media':':.2f','temp_media':':.2f','registros':True},
        size_max=50, zoom=3.2, color_continuous_scale=px.colors.sequential.thermal
    )
    fig_map.update_layout(mapbox_style='carto-positron', margin={'r':0,'t':0,'l':0,'b':0}, paper_bgcolor="rgba(0,0,0,0)")
    st.plotly_chart(fig_map, use_container_width=True)

with right:
    st.subheader("Mapa 2 — Rotas entre Capitais (linha = rota; largura ≈ tarifa)")
    # We use pydeck for smooth arcs
    arc_layer = pdk.Layer(
        "ArcLayer",
        data=routes,
        get_source_position=["olon","olat"],
        get_target_position=["dlon","dlat"],
        get_source_color=[6,110,204],
        get_target_color=[255,107,74],
        get_width="tarifa_media",
        pickable=True,
        auto_highlight=True
    )
    view = pdk.ViewState(latitude=-14.2350, longitude=-51.9253, zoom=3.4, pitch=0)
    r_deck = pdk.Deck(layers=[arc_layer], initial_view_state=view, map_style='mapbox://styles/mapbox/light-v9')
    st.pydeck_chart(r_deck)

st.markdown("---")

# -----------------------------
# Insights — estações e correlações
st.subheader("Insights — estações, top rotas e companhias")
dff['ESTACAO'] = dff['MES'].apply(lambda m: 'Verão' if m in [12,1,2] else ('Outono' if m in [3,4,5] else ('Inverno' if m in [6,7,8] else 'Primavera')))

est = (
    dff.groupby('ESTACAO')
    .agg(tarifa_media=('TARIFA','mean'), temp_media=('TEMP_MEDIA','mean'), registros=('TARIFA','count'))
    .reset_index()
)
t1, t2, t3 = st.columns(3)
with t1:
    st.markdown("### 🔁 Tarifas por estação")
    st.table(est.round(2).set_index('ESTACAO'))
with t2:
    st.markdown("### 🚩 Top rotas (freq)")
    top_routes = dff['ROTA'].value_counts().head(8).rename_axis('ROTA').reset_index(name='FREQ')
    st.table(top_routes)
with t3:
    st.markdown("### ✈️ Companhias — média e registros")
    comp = dff.groupby('COMPANHIA').agg(tarifa_media=('TARIFA','mean'), regs=('TARIFA','count')).reset_index().sort_values('tarifa_media', ascending=False)
    st.dataframe(comp.head(12).round(2))

st.markdown("---")

# -----------------------------
# Forecasting — Prophet (model per rota or per capital)
st.header("Previsão — Tarifas para 2026")
st.markdown("Escolha uma rota (ou selecione 'AGREGAR POR DESTINO' para modelos por capital). Usamos Prophet com regressor de temperatura para mais precisão.")

agg_mode = st.radio("Modo de forecast:", ("Por Rota", "Por Destino (agregado)"), index=0)
if agg_mode == "Por Rota":
    rota_sel = st.selectbox("Selecione rota", sorted(routes['ROTA'].unique()))
    df_model = dff[dff['ROTA']==rota_sel].groupby('DATA').agg(tarifa=('TARIFA','mean'), temp=('TEMP_MEDIA','mean')).reset_index()
else:
    dest_sel = st.selectbox("Selecione destino (capital)", sorted(agg_cap['CAPITAL'].unique()))
    df_model = dff[dff['DEST_ROTA']==dest_sel].groupby('DATA').agg(tarifa=('TARIFA','mean'), temp=('TEMP_MEDIA','mean')).reset_index()

if df_model.shape[0] < 12:
    st.warning("Dados insuficientes (menos de 12 meses). Escolha outra rota/destino ou aumente o período de filtro.")
else:
    df_prop = df_model.rename(columns={'DATA':'ds','tarifa':'y','temp':'temp_reg'}).dropna(subset=['ds','y'])
    m = Prophet(yearly_seasonality=True, weekly_seasonality=False, daily_seasonality=False)
    m.add_regressor('temp_reg')
    with st.spinner("Treinando modelo Prophet (poucos segundos)…"):
        m.fit(df_prop)
    future = m.make_future_dataframe(periods=12, freq='MS')
    # fill future regressor by monthly averages from history (seasonality)
    df_prop['month'] = df_prop['ds'].dt.month
    monthly_temp = df_prop.groupby('month')['temp_reg'].mean().to_dict()
    future['month'] = future['ds'].dt.month
    future['temp_reg'] = future['month'].map(monthly_temp).fillna(df_prop['temp_reg'].mean())
    forecast = m.predict(future)
    # Plot interactive
    st.plotly_chart(plot_plotly(m, forecast), use_container_width=True)
    # show 2026 table
    f2026 = forecast[forecast['ds'].dt.year==2026][['ds','yhat','yhat_lower','yhat_upper']].rename(columns={'ds':'DATA','yhat':'TARIFA_PRED'}).round(2)
    st.subheader("Previsão mês a mês 2026")
    st.table(f2026.set_index('DATA'))
    # save
    safe_name = (rota_sel if agg_mode=="Por Rota" else dest_sel).replace(" ","_").replace("/","_")
    out_csv = os.path.join(OUTPUT_PATH, f"forecast_2026_{safe_name}.csv")
    f2026.to_csv(out_csv, index=False)
    st.success(f"Previsão salva em: {out_csv}")

st.markdown("---")

# -----------------------------
# Exports: CSV + PPTX quick (speaker notes minimal)
exp1, exp2 = st.columns(2)
with exp1:
    if st.button("Exportar CSV (filtrado)"):
        outf = os.path.join(OUTPUT_PATH, "boraali_dataset_filtrado.csv")
        dff.to_csv(outf, index=False)
        st.success(f"CSV salvo: {outf}")
with exp2:
    if st.button("Gerar PPTX (4 slides)"):
        try:
            from pptx import Presentation
            prs = Presentation()
            s0 = prs.slides.add_slide(prs.slide_layouts[0])
            s0.shapes.title.text = "Bora Alí — SR2"
            s0.placeholders[1].text = "Tarifas, clima e previsão — Capitais"
            s1 = prs.slides.add_slide(prs.slide_layouts[1]); s1.shapes.title.text = "Dados & Metodologia"
            s1.placeholders[1].text = "Integração ANAC + INMET. Limpeza: outliers e padronização."
            s2 = prs.slides.add_slide(prs.slide_layouts[1]); s2.shapes.title.text = "Insights"
            s2.placeholders[1].text = "Estações vs tarifas; Top rotas; Companhias"
            s3 = prs.slides.add_slide(prs.slide_layouts[1]); s3.shapes.title.text = "Previsão 2026"
            s3.placeholders[1].text = "Prophet (regressor: temperatura)."
            pptx_path = os.path.join(OUTPUT_PATH, "BoraAli_SR2_slides.pptx")
            prs.save(pptx_path)
            st.success(f"PPTX salvo: {pptx_path}")
        except Exception as e:
            st.error(f"Erro ao gerar PPTX: {e}")

st.markdown("**Design:** linguagem 'Bora Alí' — legendas descoladas e diretas; foco em clareza e storytelling.")
st.markdown("Made with 💙 by Bora Alí — boa sorte na SR2 ✈️")

